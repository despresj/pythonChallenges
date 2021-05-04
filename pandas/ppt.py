import pandas as pd
import io
import requests
from pptx import Presentation
from datetime import datetime, timedelta

def yesterday(outformat=False):
    yesterday = datetime.now() - timedelta(1)
    if outformat:
        return yesterday.strftime('%m-%d-%Y')
    else:
        return yesterday.strftime('%Y-%m-%d')

yesterday()
yesterday(outformat=True)

url = "https://raw.githubusercontent.com/nytimes/covid-19-data/master/us-states.csv"
s = requests.get(url).content
df = pd.read_csv(io.StringIO(s.decode('utf-8')))

mi = df[(df['date'] == yesterday()) & (df['state'] == 'Michigan')]
deaths = str(mi['deaths'].values[0])

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = yesterday(outformat=True)

tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = "lets see"
p.level = 1

p = tf.add_paragraph()
p.text = 'Unfortunately we lost ' + deaths + ' people to covid'
p.level = 2

file_path = '/Users/josephdespres/Documents/Python/pythonChallenges/pandas/ppt/' + yesterday() + '_covid.pptx'
prs.save(file_path)